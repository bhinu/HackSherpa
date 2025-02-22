import os
from pptx import Presentation  # For creating PowerPoint presentations
from pptx.util import Inches  # For positioning elements
import subprocess


def create_presentation(readme_content, template_path=None):
    """
    Create a PowerPoint presentation from the README content.
    """
    # Load the template if provided
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()  # Create a blank presentation

    # Add a title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project Presentation"
    subtitle.text = "Generated from README.md"

    # Split README content into sections
    sections = readme_content.split("\n## ")  # Split by Markdown headings
    for section in sections:
        if not section.strip():
            continue

        # Add a new slide for each section
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide layout
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set the title and content
        title.text = section.split("\n")[0].strip()  # First line as title
        content.text = "\n".join(section.split("\n")[1:]).strip()  # Rest as content

    return prs

def save_presentation(prs, output_path):
    """
    Save the PowerPoint presentation to a file.
    """
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """
    Convert a PowerPoint presentation to PDF using LibreOffice.
    """
    try:
        # Use LibreOffice to convert PPTX to PDF
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", os.path.dirname(pdf_path)], check=True)
        print(f"Presentation converted to PDF: {pdf_path}")
    except subprocess.CalledProcessError as e:
        print(f"Error: Failed to convert PPTX to PDF. {e}")

def generate_presentation_from_readme(readme_content, output_dir, template_path=None):
    """
    Generate a PowerPoint presentation and PDF from README content.
    """
    # Create the presentation
    prs = create_presentation(readme_content, template_path)

    # Save the presentation as PPTX
    pptx_path = os.path.join(output_dir, "presentation.pptx")
    save_presentation(prs, pptx_path)

    # Convert the presentation to PDF
    pdf_path = os.path.join(output_dir, "presentation.pdf")
    convert_pptx_to_pdf(pptx_path, pdf_path)